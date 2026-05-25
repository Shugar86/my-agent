#!/usr/bin/env python3
"""Create PPTX presentation about ABS vs PLA vs PETG using python-pptx directly."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color theme
PRIMARY = RGBColor(31, 78, 120)      # #1F4E78
SECONDARY = RGBColor(46, 117, 182)    # #2E75B6
ACCENT = RGBColor(197, 90, 17)        # #C55A11
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(31, 41, 55)      # #1f2937
GRAY_TEXT = RGBColor(107, 114, 128)   # #6b7280

# PLA green, ABS orange, PETG blue
PLA_COLOR = RGBColor(198, 239, 206)   # #C6EFCE
ABS_COLOR = RGBColor(255, 235, 156)   # #FFEB9C
PETG_COLOR = RGBColor(189, 215, 238)  # #BDD7EE


def add_title_slide(prs, title, subtitle, footer=""):
    """Add a title slide with gradient-like dark blue background."""
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 240)
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = foot_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(180, 200, 220)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, lines, footer=""):
    """Add a content slide with bullet points."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
        
        # Bold lines starting with • or containing "Вывод:"
        if line.startswith('•') or 'Вывод:' in line or 'Важно:' in line:
            p.font.bold = True
        if 'Вывод:' in line or 'Важно:' in line:
            p.font.color.rgb = ACCENT
    
    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4))
        tf = foot_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


def add_table_slide(prs, title, headers, rows, footer=""):
    """Add a slide with a comparison table."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    
    table = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.4), Inches(12.333), Inches(4.8)).table
    
    # Set column widths
    col_width = 12.333 / cols_count
    for col in table.columns:
        col.width = Inches(col_width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Color code by material column
            if col_idx == 1 and 'PLA' in str(headers[col_idx]):
                cell.fill.solid()
                cell.fill.fore_color.rgb = PLA_COLOR
            elif col_idx == 2 and 'ABS' in str(headers[col_idx]):
                cell.fill.solid()
                cell.fill.fore_color.rgb = ABS_COLOR
            elif col_idx == 3 and 'PETG' in str(headers[col_idx]):
                cell.fill.solid()
                cell.fill.fore_color.rgb = PETG_COLOR
    
    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4))
        tf = foot_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


def add_three_column_slide(prs, title, col1_title, col1_items, col2_title, col2_items, col3_title, col3_items):
    """Add a slide with three columns."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Three columns
    columns = [
        (col1_title, col1_items, PLA_COLOR, RGBColor(55, 86, 35)),
        (col2_title, col2_items, ABS_COLOR, RGBColor(132, 60, 12)),
        (col3_title, col3_items, PETG_COLOR, RGBColor(31, 78, 120)),
    ]
    
    col_width = 3.8
    for i, (col_title, items, bg_color, text_color) in enumerate(columns):
        x = 0.5 + i * (col_width + 0.3)
        
        # Background box
        box = slide.shapes.add_shape(1, Inches(x), Inches(1.4), Inches(col_width), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.5), Inches(col_width - 0.2), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = text_color
        
        # Items
        items_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.1), Inches(col_width - 0.2), Inches(4.0))
        tf = items_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(6)
    
    return slide


def add_matrix_slide(prs, title, headers, rows):
    """Add a star-rating matrix slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Matrix table
    rows_count = len(rows) + 1
    cols_count = len(headers)
    
    table = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.0)).table
    
    col_width = 12.333 / cols_count
    for col in table.columns:
        col.width = Inches(col_width)
    
    # Header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
            
            # Color code
            if col_idx == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PLA_COLOR
            elif col_idx == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ABS_COLOR
            elif col_idx == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PETG_COLOR
    
    return slide


# ========== CREATE PRESENTATION ==========

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, 
    "Сравнение материалов для 3D печати", 
    "ABS vs PLA vs PETG",
    "Актуально на май 2025 | Данные: Ozon, Wildberries, 3DTool, 3DTech | My Agent AI")

# Slide 2: Overview Table
add_table_slide(prs,
    "Обзор материалов",
    ["Параметр", "PLA (Полилактид)", "ABS", "PETG (ПЭТГ)"],
    [
        ["Плотность", "1.24 г/см³", "1.04 г/см³", "1.27 г/см³"],
        ["Темп. стеклования", "55-60°C", "105°C", "75-80°C"],
        ["Темп. печати", "190-220°C", "220-250°C", "220-250°C"],
        ["Темп. стола", "20-60°C (опц.)", "90-110°C (обяз.)", "70-80°C (реком.)"],
    ],
    "My Agent AI | NeuroAPI + Tavily")

# Slide 3: Mechanical Properties
add_table_slide(prs,
    "Механические характеристики",
    ["Параметр", "PLA", "ABS", "PETG"],
    [
        ["Прочность на растяжение", "50-60 МПа", "40-50 МПа", "50-60 МПа"],
        ["Модуль упругости", "2.5-3.5 ГПа", "2.0-2.5 ГПа", "2.0-2.8 ГПа"],
        ["Удлинение при разрыве", "3-5%", "10-25%", "15-30%"],
        ["Влагопоглощение", "0.5%", "0.3%", "0.2%"],
    ],
    "Источник: 3DToday, 3DTool technical datasheets")

# Slide 4: Safety
add_content_slide(prs,
    "Безопасность и экология",
    [
        "Токсичность при печати:",
        "• PLA: низкая, безопасен",
        "• ABS: высокая, выделяет пары (нужна вентиляция!)",
        "• PETG: низкая, безопасен",
        "",
        "Запах:",
        "• PLA: приятный (попкорн)",
        "• ABS: резкий (пластик)",
        "• PETG: почти нет",
        "",
        "Пищевая безопасность: PLA и PETG — да, ABS — нет",
        "Биоразлагаемость: только PLA",
        "УФ-стойкость: ABS лучший, PETG средний, PLA низкий",
    ])

# Slide 5: Prices
add_table_slide(prs,
    "Цены в России (2024-2025)",
    ["Магазин / Бренд", "PLA", "ABS", "PETG"],
    [
        ["Ozon (ESUN/FDplast)", "750 ₽/кг", "900 ₽/кг", "850 ₽/кг"],
        ["Ozon (Bestfilament)", "950 ₽/кг", "1200 ₽/кг", "1150 ₽/кг"],
        ["Wildberries (NoName)", "650 ₽/кг", "800 ₽/кг", "750 ₽/кг"],
        ["3DTool Premium", "—", "1700 ₽/кг", "1400 ₽/кг"],
        ["Премиум (ColorFabb)", "2200 ₽/кг", "2500 ₽/кг", "2400 ₽/кг"],
    ],
    "Источник: Ozon, Wildberries, 3DTool, 3DTech (май 2025)")

# Slide 6: Use Cases (Three columns)
add_three_column_slide(prs,
    "Рекомендации по применению",
    "PLA",
    ["Прототипы и макеты", "Декоративные изделия", "Контейнеры для еды", "Начинающим пользователям", "Biodegradable проекты"],
    "ABS",
    ["Автомобильные запчасти", "Функциональные детали", "Уличные конструкции", "Детали на солнце (УФ)", "Термостойкие элементы"],
    "PETG",
    ["Детали для кухни/ванной", "Влагостойкие элементы", "Механически нагруженные", "Золотая середина!", "Пищевая безопасность"]
)

# Slide 7: Matrix
add_matrix_slide(prs,
    "Сводная матрица выбора",
    ["Критерий", "PLA", "ABS", "PETG"],
    [
        ["Простота печати", "★★★★★", "★★☆☆☆", "★★★★☆"],
        ["Прочность", "★★★☆☆", "★★★★☆", "★★★★★"],
        ["Термостойкость", "★★☆☆☆", "★★★★★", "★★★☆☆"],
        ["Безопасность", "★★★★★", "★★☆☆☆", "★★★★★"],
        ["Ударостойкость", "★★☆☆☆", "★★★★☆", "★★★★★"],
        ["Влагостойкость", "★★☆☆☆", "★★★☆☆", "★★★★★"],
        ["Цена", "★★★★★", "★★★☆☆", "★★★★☆"],
    ])

# Slide 8: Conclusion
add_title_slide(prs,
    "Итоги и выводы",
    "",
    "")

# Add conclusion text to slide 8 (it's the last slide)
last_slide = prs.slides[-1]

# Three conclusion boxes
conclusions = [
    ("PLA", "Лучший для начинающих и декора. Безопасен, дешев, биоразлагаем. Не подходит для нагрева выше 60°C.", PLA_COLOR),
    ("ABS", "Для профи: термостойкость, УФ-защита, прочность. Требует камеры с вентиляцией и нагреваемого стола.", ABS_COLOR),
    ("PETG", "Золотая середина! Прочность ABS + безопасность PLA. Идеален для кухни, ванной, функциональных деталей.", PETG_COLOR),
]

box_width = 3.8
for i, (title, text, bg_color) in enumerate(conclusions):
    x = 0.5 + i * (box_width + 0.3)
    
    # Background
    box = last_slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(box_width), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()
    
    # Title
    title_box = last_slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.9), Inches(box_width - 0.2), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    
    # Text
    text_box = last_slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.5), Inches(box_width - 0.2), Inches(2.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT

# Footer on last slide
foot_box = last_slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
tf = foot_box.text_frame
p = tf.paragraphs[0]
p.text = "Создано с помощью My Agent AI | API: NeuroAPI + OpenRouter + Tavily | Май 2025"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(180, 200, 220)
p.alignment = PP_ALIGN.CENTER

# Save
output_path = 'output/презентация_ABS_PLA_PETG.pptx'
prs.save(output_path)

print(f'✅ Презентация создана: {output_path}')
print(f'📊 Размер: {os.path.getsize(output_path):,} bytes ({os.path.getsize(output_path)/1024:.1f} KB)')
print(f'🖼️  Слайдов: {len(prs.slides)}')
print(f'📍 Путь: {os.path.abspath(output_path)}')

import os
import json
from typing import List, Dict, Any, Optional
from pathlib import Path


def create_slide_deck(title: str, slides: List[Dict[str, Any]], 
                      theme: Dict[str, str] = None) -> Dict[str, Any]:
    """Create HTML slide deck structure."""
    
    default_theme = {
        "primary": "#2563eb",
        "secondary": "#7c3aed", 
        "background": "#ffffff",
        "text": "#1f2937",
        "accent": "#f59e0b",
        "font_heading": "Space Grotesk, sans-serif",
        "font_body": "Inter, sans-serif"
    }
    theme = theme or default_theme
    
    # Generate HTML for each slide
    html_slides = []
    for i, slide in enumerate(slides):
        slide_html = _render_slide(slide, theme, i + 1, len(slides))
        html_slides.append({
            "number": i + 1,
            "type": slide.get("type", "content"),
            "html": slide_html
        })
    
    # Build full deck HTML
    deck_html = _build_deck_html(title, html_slides, theme)
    
    return {
        "title": title,
        "slide_count": len(slides),
        "theme": theme,
        "slides": html_slides,
        "html": deck_html
    }


def _render_slide(slide: Dict[str, Any], theme: Dict[str, str], 
                  number: int, total: int) -> str:
    """Render individual slide HTML."""
    slide_type = slide.get("type", "content")
    
    if slide_type == "title":
        return f"""
        <div class="slide title-slide" data-type="title">
            <div class="slide-content">
                <h1 style="font-family: {theme['font_heading']}; color: {theme['primary']}; font-size: 48px; margin-bottom: 20px;">
                    {slide.get('title', '')}
                </h1>
                <p style="font-family: {theme['font_body']}; color: {theme['text']}; font-size: 24px; opacity: 0.8;">
                    {slide.get('subtitle', '')}
                </p>
            </div>
        </div>
        """
    
    elif slide_type == "content":
        bullets = slide.get("bullets", [])
        bullets_html = "\n".join([f'<li style="margin-bottom: 12px; font-size: 20px;">{b}</li>' for b in bullets])
        
        return f"""
        <div class="slide content-slide" data-type="content">
            <div class="slide-content">
                <h2 style="font-family: {theme['font_heading']}; color: {theme['primary']}; font-size: 36px; margin-bottom: 30px;">
                    {slide.get('title', '')}
                </h2>
                <ul style="font-family: {theme['font_body']}; color: {theme['text']}; line-height: 1.6;">
                    {bullets_html}
                </ul>
            </div>
        </div>
        """
    
    elif slide_type == "chart":
        chart_data = slide.get("data", {})
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        
        # Simple bar chart using divs
        max_val = max(values) if values else 1
        bars_html = ""
        for label, value in zip(labels, values):
            width = (value / max_val) * 100 if max_val > 0 else 0
            bars_html += f'''
            <div style="margin-bottom: 15px;">
                <div style="display: flex; justify-content: space-between; margin-bottom: 5px;">
                    <span style="font-family: {theme['font_body']}; font-size: 14px;">{label}</span>
                    <span style="font-family: {theme['font_body']}; font-size: 14px; font-weight: bold;">{value}</span>
                </div>
                <div style="background: #e5e7eb; border-radius: 4px; height: 30px;">
                    <div style="background: {theme['primary']}; width: {width}%; height: 100%; border-radius: 4px; transition: width 0.3s;"></div>
                </div>
            </div>
            '''
        
        return f"""
        <div class="slide chart-slide" data-type="chart">
            <div class="slide-content">
                <h2 style="font-family: {theme['font_heading']}; color: {theme['primary']}; font-size: 36px; margin-bottom: 30px;">
                    {slide.get('title', '')}
                </h2>
                <div style="font-family: {theme['font_body']}; padding: 20px;">
                    {bars_html}
                </div>
            </div>
        </div>
        """
    
    elif slide_type == "two_column":
        left = slide.get("left", {})
        right = slide.get("right", {})
        
        return f"""
        <div class="slide two-column-slide" data-type="two_column">
            <div class="slide-content" style="display: flex; gap: 40px;">
                <div style="flex: 1;">
                    <h3 style="font-family: {theme['font_heading']}; color: {theme['primary']}; font-size: 28px; margin-bottom: 20px;">
                        {left.get('title', '')}
                    </h3>
                    <p style="font-family: {theme['font_body']}; color: {theme['text']}; font-size: 18px; line-height: 1.6;">
                        {left.get('content', '')}
                    </p>
                </div>
                <div style="flex: 1;">
                    <h3 style="font-family: {theme['font_heading']}; color: {theme['secondary']}; font-size: 28px; margin-bottom: 20px;">
                        {right.get('title', '')}
                    </h3>
                    <p style="font-family: {theme['font_body']}; color: {theme['text']}; font-size: 18px; line-height: 1.6;">
                        {right.get('content', '')}
                    </p>
                </div>
            </div>
        </div>
        """
    
    else:
        return f"""
        <div class="slide" data-type="{slide_type}">
            <div class="slide-content">
                <h2 style="font-family: {theme['font_heading']}; color: {theme['primary']}; font-size: 36px;">
                    {slide.get('title', '')}
                </h2>
                <p style="font-family: {theme['font_body']}; color: {theme['text']}; font-size: 18px;">
                    {slide.get('content', '')}
                </p>
            </div>
        </div>
        """


def _build_deck_html(title: str, slides: List[Dict[str, Any]], 
                     theme: Dict[str, str]) -> str:
    """Build complete deck HTML."""
    
    slides_html = "\n".join([s["html"] for s in slides])
    
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link href="https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;600;700&family=Inter:wght@400;500;600&display=swap" rel="stylesheet">
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ 
            font-family: {theme['font_body']}; 
            background: #f3f4f6;
            padding: 20px;
        }}
        .slide {{
            width: 960px;
            height: 540px;
            margin: 0 auto 30px;
            background: {theme['background']};
            border-radius: 8px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            padding: 60px;
            display: flex;
            align-items: center;
            justify-content: center;
            position: relative;
            overflow: hidden;
        }}
        .slide::before {{
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            height: 6px;
            background: {theme['primary']};
        }}
        .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 30px;
            font-size: 14px;
            color: #9ca3af;
        }}
        .title-slide {{
            text-align: center;
            background: {theme['primary']};
        }}
        .title-slide::before {{
            background: {theme['accent']};
        }}
        .title-slide h1 {{
            color: white !important;
        }}
        .title-slide p {{
            color: rgba(255,255,255,0.9) !important;
        }}
        @media print {{
            .slide {{
                page-break-after: always;
                margin: 0;
                box-shadow: none;
            }}
        }}
    </style>
</head>
<body>
    {slides_html}
</body>
</html>"""


def export_to_pptx(deck: Dict[str, Any], output_path: str = None) -> str:
    """Export slide deck to PPTX format."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        theme = deck.get("theme", {})
        primary = theme.get("primary", "#2563eb")
        
        for slide_data in deck.get("slides", []):
            slide_type = slide_data.get("type", "content")
            
            # Add blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
            
            # Add title bar
            title_bar = slide.shapes.add_shape(
                1, Inches(0), Inches(0), Inches(13.333), Inches(0.1)
            )
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = RGBColor(
                int(primary[1:3], 16), int(primary[3:5], 16), int(primary[5:7], 16)
            )
            
            if slide_type == "title":
                _add_title_slide(slide, slide_data, theme)
            elif slide_type == "content":
                _add_content_slide(slide, slide_data, theme)
            else:
                _add_generic_slide(slide, slide_data, theme)
        
        if not output_path:
            output_path = f"output/{deck['title'].replace(' ', '_')}.pptx"
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        
        prs.save(output_path)
        return output_path
        
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"Error exporting PPTX: {str(e)}"


def _add_title_slide(slide, slide_data, theme):
    """Add title slide to PPTX."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("html", "").split("<h1>")[1].split("</h1>")[0] if "<h1>" in slide_data.get("html", "") else "Title"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = 1  # Center
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    subtitle = slide_data.get("html", "").split("<p>")[1].split("</p>")[0] if "<p>" in slide_data.get("html", "") else ""
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = 1  # Center


def _add_content_slide(slide, slide_data, theme):
    """Add content slide to PPTX."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    # Use explicit title/content keys if available, fallback to HTML parsing
    title = slide_data.get("title", "")
    content = slide_data.get("content", "")
    
    # Fallback: extract from HTML if keys not present
    if not title and not content:
        html = slide_data.get("html", "")
        if "<h2>" in html:
            title = html.split("<h2>")[1].split("</h2>")[0]
            title = title.split(">")[-1].split("<")[0] if ">" in title else title
        if "<ul>" in html:
            import re
            content = re.sub('<[^<]+?>', '', html)
            content = content.strip()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.strip()
    p.font.size = Pt(32)
    p.font.bold = True
    primary = theme.get("primary", "#2563eb")
    p.font.color.rgb = RGBColor(int(primary[1:3], 16), int(primary[3:5], 16), int(primary[5:7], 16))
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content[:1000]  # Increased limit
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(31, 41, 55)


def _add_generic_slide(slide, slide_data, theme):
    """Add generic slide to PPTX."""
    _add_content_slide(slide, slide_data, theme)


def save_slide_html(deck: Dict[str, Any], output_dir: str = "output/slides") -> str:
    """Save slide deck as HTML files."""
    os.makedirs(output_dir, exist_ok=True)
    
    # Save full deck
    deck_path = os.path.join(output_dir, "deck.html")
    with open(deck_path, "w", encoding="utf-8") as f:
        f.write(deck["html"])
    
    # Save individual slides
    for slide in deck["slides"]:
        slide_path = os.path.join(output_dir, f"slide_{slide['number']:02d}.html")
        with open(slide_path, "w", encoding="utf-8") as f:
            f.write(slide["html"])
    
    return output_dir


def register_tools():
    from tools import slides_tools
    slides_tools.register_tools()


def unregister_tools():
    from tools import slides_tools
    slides_tools.unregister_tools()
